import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor

# Buraya hutbe gelecek
metin = """
Duyarlılık önemli bir erdemdir.
Müslüman sorumluluk sahibidir.
Allah iyilik yapanları sever.
"""

# Cümlelere ayır
cumleler = re.split(r'(?<=[.!?])\s+', metin.strip())

sunum = Presentation()

for cumle in cumleler:

    slayt = sunum.slides.add_slide(sunum.slide_layouts[6])

    # Üst mavi şerit
    serit = slayt.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        sunum.slide_width,
        500000
    )

    serit.fill.solid()
    serit.fill.fore_color.rgb = RGBColor(31, 78, 121)

    kutu = slayt.shapes.add_textbox(
        Inches(1),
        Inches(1.3),
        Inches(11),
        Inches(4)
    )

    tf = kutu.text_frame
    tf.text = cumle

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    p.font.size = Pt(30)

sunum.save("HutbeSunumu.pptx")

print("Bitti.")