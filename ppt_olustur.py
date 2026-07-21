import re
import sys
from datetime import datetime
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ayarlar import YAZI_TIPI


KORUNAN_IFADELER = [
    "s.a.s.",
    "s.a.s",
    "s.a.v.",
    "s.a.v",
    "a.s.",
    "a.s",
    "r.a.",
    "r.a",
    "Hz.",
    "hz.",
    "Dr.",
    "Prof.",
    "vb.",
    "vs.",
]


def arapca_harf_mi(karakter):
    return (
        "\u0600" <= karakter <= "\u06ff"
        or "\u0750" <= karakter <= "\u077f"
        or "\u08a0" <= karakter <= "\u08ff"
        or "\ufb50" <= karakter <= "\ufdff"
        or "\ufe70" <= karakter <= "\ufeff"
    )


def arapca_iceriyor_mu(metin):
    return any(arapca_harf_mi(karakter) for karakter in metin)


def hitap_mi(metin):
    hitaplar = (
        "Aziz Müminler!",
        "Aziz Müslümanlar!",
        "Kıymetli Müslümanlar!",
        "Kıymetli Müminler!",
        "Kıymetli Kardeşlerim!",
        "Değerli Müminler!",
        "Değerli Müslümanlar!",
        "Değerli Kardeşlerim!",
        "Muhterem Müslümanlar!",
    )

    return metin.strip() in hitaplar


def kucuk_hatalari_duzelt(metin):
    duzeltmeler = {
        "(s.a.s._2)": "(s.a.s.)",
        "s.a.s._2": "s.a.s.",
        "(s.a.v._2)": "(s.a.v.)",
        "s.a.v._2": "s.a.v.",
        "(a.s._2)": "(a.s.)",
        "a.s._2": "a.s.",
        "(r.a._2)": "(r.a.)",
        "r.a._2": "r.a.",
    }

    sonuc = metin

    for hatali, dogru in duzeltmeler.items():
        sonuc = sonuc.replace(hatali, dogru)

    return sonuc


def korunan_ifadeleri_degistir(metin):
    sonuc = metin
    eslestirmeler = {}

    ifadeler = sorted(
        KORUNAN_IFADELER,
        key=len,
        reverse=True,
    )

    for sira, ifade in enumerate(ifadeler):
        kod = f"KORUMA{sira:03d}X"
        sonuc = sonuc.replace(ifade, kod)
        eslestirmeler[kod] = ifade

    return sonuc, eslestirmeler


def korunan_ifadeleri_geri_getir(metin, eslestirmeler):
    sonuc = metin

    for kod, ifade in eslestirmeler.items():
        sonuc = sonuc.replace(kod, ifade)

    return sonuc


def metni_temizle(metin):
    metin = kucuk_hatalari_duzelt(metin)
    metin = metin.replace("\r\n", "\n").replace("\r", "\n")
    metin = re.sub(r"[ \t]+", " ", metin)

    temiz_satirlar = []

    for satir in metin.split("\n"):
        satir = satir.strip()

        if not satir:
            continue

        if arapca_iceriyor_mu(satir):
            continue

        if re.fullmatch(r"\d+", satir):
            continue

        temiz_satirlar.append(satir)

    sonuc = " ".join(temiz_satirlar)
    return kucuk_hatalari_duzelt(sonuc)


def cumlelere_ayir(metin):
    temiz_metin = metni_temizle(metin)

    if not temiz_metin:
        return []

    korunmus_metin, eslestirmeler = korunan_ifadeleri_degistir(
        temiz_metin
    )

    parcalar = re.split(
        r'(?<=[.!?])(?:["”’\']*)\s+',
        korunmus_metin,
    )

    sonuc = []

    for parca in parcalar:
        parca = korunan_ifadeleri_geri_getir(
            parca,
            eslestirmeler,
        )

        parca = kucuk_hatalari_duzelt(parca)
        parca = " ".join(parca.split()).strip()

        if not parca:
            continue

        if arapca_iceriyor_mu(parca):
            continue

        if re.fullmatch(r"\d+", parca):
            continue

        sonuc.append(parca)

    return sonuc


def yazi_boyutu_belirle(metin):
    uzunluk = len(metin)

    if hitap_mi(metin):
        return 46

    if uzunluk <= 35:
        return 46

    if uzunluk <= 60:
        return 42

    if uzunluk <= 90:
        return 38

    if uzunluk <= 125:
        return 34

    if uzunluk <= 170:
        return 30

    if uzunluk <= 220:
        return 27

    return 24


def placeholderlari_sil(slayt):
    silinecekler = [
        sekil
        for sekil in slayt.shapes
        if sekil.is_placeholder
    ]

    for sekil in silinecekler:
        element = sekil._element
        element.getparent().remove(element)


def slayttaki_tum_sekilleri_sil(slayt):
    for sekil in list(slayt.shapes):
        element = sekil._element
        element.getparent().remove(element)


def slayta_metin_ekle(slayt, metin):
    metin = kucuk_hatalari_duzelt(metin)

    kutu = slayt.shapes.add_textbox(
        Inches(0.55),
        Inches(1.25),
        Inches(8.90),
        Inches(5.25),
    )

    metin_alani = kutu.text_frame
    metin_alani.clear()
    metin_alani.word_wrap = True
    metin_alani.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    metin_alani.margin_left = Inches(0.08)
    metin_alani.margin_right = Inches(0.08)
    metin_alani.margin_top = Inches(0.05)
    metin_alani.margin_bottom = Inches(0.05)

    paragraf = metin_alani.paragraphs[0]
    paragraf.text = metin
    paragraf.alignment = PP_ALIGN.CENTER
    paragraf.level = 0

    boyut = yazi_boyutu_belirle(metin)

    for run in paragraf.runs:
        run.font.name = YAZI_TIPI
        run.font.size = Pt(boyut)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)


def sayisal_dosya_sirasi(dosya_yolu):
    eslesme = re.search(r"\d+", dosya_yolu.stem)

    if eslesme:
        return (
            0,
            int(eslesme.group()),
            dosya_yolu.name.lower(),
        )

    return (
        1,
        0,
        dosya_yolu.name.lower(),
    )


def klasordeki_gorselleri_bul(
    klasor_yolu,
    izinli_uzantilar,
):
    if not klasor_yolu.exists():
        return []

    gorseller = [
        dosya
        for dosya in klasor_yolu.iterdir()
        if (
            dosya.is_file()
            and dosya.suffix.lower() in izinli_uzantilar
        )
    ]

    return sorted(
        gorseller,
        key=sayisal_dosya_sirasi,
    )


def slayta_tam_gorsel_ekle(
    slayt,
    gorsel_yolu,
    slayt_genisligi,
    slayt_yuksekligi,
):
    slayttaki_tum_sekilleri_sil(slayt)

    with Image.open(gorsel_yolu) as resim:
        resim_genisligi, resim_yuksekligi = resim.size

    resim_orani = resim_genisligi / resim_yuksekligi
    slayt_orani = slayt_genisligi / slayt_yuksekligi

    if resim_orani > slayt_orani:
        yukseklik = slayt_yuksekligi
        genislik = int(yukseklik * resim_orani)
        sol = int((slayt_genisligi - genislik) / 2)
        ust = 0

    else:
        genislik = slayt_genisligi
        yukseklik = int(genislik / resim_orani)
        sol = 0
        ust = int((slayt_yuksekligi - yukseklik) / 2)

    slayt.shapes.add_picture(
        str(gorsel_yolu),
        sol,
        ust,
        width=genislik,
        height=yukseklik,
    )


def uygulama_klasorunu_bul():
    """
    Python ile çalışırken proje klasörünü,
    EXE ile çalışırken arayuz.exe'nin bulunduğu klasörü döndürür.
    """
    if getattr(sys, "frozen", False):
        return Path(sys.executable).resolve().parent

    return Path(__file__).resolve().parent


def masaustu_klasorunu_bul():
    ev = Path.home()

    olasi_yollar = [
        ev / "Desktop",
        ev / "Masaüstü",
        ev / "OneDrive" / "Desktop",
        ev / "OneDrive" / "Masaüstü",
    ]

    for yol in olasi_yollar:
        if yol.exists():
            return yol

    return ev


def benzersiz_dosya_yolu_olustur():
    tarih = datetime.now().strftime("%d-%m-%Y")
    masaustu = masaustu_klasorunu_bul()

    temel_ad = f"{tarih}_HutbeSunumu"
    dosya_yolu = masaustu / f"{temel_ad}.pptx"

    sayac = 2

    while dosya_yolu.exists():
        dosya_yolu = masaustu / f"{temel_ad}_{sayac}.pptx"
        sayac += 1

    return dosya_yolu


def sunum_olustur(metin):
    proje_klasoru = uygulama_klasorunu_bul()

    sablon_yolu = proje_klasoru / "hutbe_sablonu.pptx"
    giris_klasoru = proje_klasoru / "giris_gorselleri"
    cikis_klasoru = proje_klasoru / "cikis_gorselleri"

    if not sablon_yolu.exists():
        raise FileNotFoundError(
            "hutbe_sablonu.pptx dosyası "
            "uygulama klasöründe bulunamadı."
        )

    cumleler = cumlelere_ayir(metin)

    if not cumleler:
        raise ValueError(
            "Sunuma eklenecek Türkçe metin bulunamadı."
        )

    giris_gorselleri = klasordeki_gorselleri_bul(
        giris_klasoru,
        {".png"},
    )

    cikis_gorselleri = klasordeki_gorselleri_bul(
        cikis_klasoru,
        {".jpg", ".jpeg"},
    )

    sunum = Presentation(sablon_yolu)

    if len(sunum.slides) == 0:
        raise ValueError(
            "Şablonun içinde en az bir slayt bulunmalıdır."
        )

    sablon_duzeni = sunum.slides[0].slide_layout
    ilk_slayt = sunum.slides[0]

    if giris_gorselleri:
        slayta_tam_gorsel_ekle(
            ilk_slayt,
            giris_gorselleri[0],
            sunum.slide_width,
            sunum.slide_height,
        )

        for gorsel_yolu in giris_gorselleri[1:]:
            slayt = sunum.slides.add_slide(sablon_duzeni)

            slayta_tam_gorsel_ekle(
                slayt,
                gorsel_yolu,
                sunum.slide_width,
                sunum.slide_height,
            )

        for cumle in cumleler:
            slayt = sunum.slides.add_slide(sablon_duzeni)
            placeholderlari_sil(slayt)
            slayta_metin_ekle(slayt, cumle)

    else:
        placeholderlari_sil(ilk_slayt)
        slayta_metin_ekle(ilk_slayt, cumleler[0])

        for cumle in cumleler[1:]:
            slayt = sunum.slides.add_slide(sablon_duzeni)
            placeholderlari_sil(slayt)
            slayta_metin_ekle(slayt, cumle)

    for gorsel_yolu in cikis_gorselleri:
        slayt = sunum.slides.add_slide(sablon_duzeni)

        slayta_tam_gorsel_ekle(
            slayt,
            gorsel_yolu,
            sunum.slide_width,
            sunum.slide_height,
        )

    dosya_yolu = benzersiz_dosya_yolu_olustur()
    sunum.save(str(dosya_yolu))

    return str(dosya_yolu)